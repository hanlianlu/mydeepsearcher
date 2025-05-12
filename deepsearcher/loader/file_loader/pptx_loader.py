"""
deepsearcher/loader/file_loader/pptx_loader.py

AdvancedPPTXLoader
------------------
• Extracts native slide text from .pptx files.  
• Runs **EasyOCR** on embedded pictures (GPU if available, else CPU).  
• Optional in‑loader chunking (`chunk_size`, `chunk_overlap`).  
• Gracefully degrades: if EasyOCR can’t start, logs a warning and skips OCR.

Dependencies
------------
pip install python-pptx pillow easyocr
(EasyOCR pulls in torch; the CPU wheel installs automatically on non‑GPU boxes.)
"""

from __future__ import annotations

import io
import logging
from typing import List, Optional

from PIL import Image
from pptx import Presentation
from langchain_core.documents import Document

# ------------------ OCR imports (optional) ------------------ #
try:
    import torch  # required by easyocr
    from easyocr import Reader
except ImportError as _imp_err:  # OCR not available
    Reader = None          # type: ignore[assignment]
    torch = None           # type: ignore[assignment]
    _EASYOCR_ERR = _imp_err
else:
    _EASYOCR_ERR = None

# ------------------ Logging ------------------ #
logger = logging.getLogger(__name__)


class AdvancedPPTXLoader:
    """
    Parameters
    ----------
    ocr_images : bool
        Run OCR over slide images if True (default).
    languages : List[str] | None
        EasyOCR language codes. Defaults to ["en"].
    gpu : bool | None
        True  → force CUDA, False → force CPU, None → auto‑detect.
    chunk_size : int | None
        If given, split slide text into chunks of this many chars
        with `chunk_overlap` overlap. If None (default), keeps one
        Document per slide.
    chunk_overlap : int
        Overlap for the optional splitter (default 200).
    """

    def __init__(
        self,
        *,
        ocr_images: bool = True,
        languages: Optional[List[str]] = None,
        gpu: bool | None = None,
        chunk_size: int | None = None,
        chunk_overlap: int = 200,
    ):
        self.ocr_images = ocr_images and Reader is not None
        self.reader = None
        if self.ocr_images:
            lang_list = languages or ["en"]
            want_gpu = torch.cuda.is_available() if gpu is None else gpu

            # Attempt GPU then CPU (unless GPU explicitly False)
            for attempt_gpu in (want_gpu, False) if want_gpu else (False,):
                try:
                    self.reader = Reader(lang_list=lang_list, gpu=attempt_gpu, verbose=False)
                    break
                except Exception as e:  # noqa: BLE001
                    logger.warning("EasyOCR init failed on %s (%s)", "GPU" if attempt_gpu else "CPU", e)
            else:
                logger.warning("EasyOCR unavailable (%s); image OCR disabled", _EASYOCR_ERR or "unknown error")
                self.ocr_images = False

        # Optional in‑loader splitter
        if chunk_size:
            from langchain.text_splitter import RecursiveCharacterTextSplitter

            self.splitter = RecursiveCharacterTextSplitter(
                chunk_size=chunk_size,
                chunk_overlap=chunk_overlap,
                separators=["\n\n", "\n", ". ", "!", "?", " ", ""],
            )
        else:
            self.splitter = None

    # ------------------------------------------------------------------ #
    # Public loader
    # ------------------------------------------------------------------ #
    def load_pptx(self, file_path: str) -> List[Document]:
        try:
            pres = Presentation(file_path)
        except Exception as e:  # noqa: BLE001
            logger.error("Failed to open PPTX %s: %s", file_path, e)
            return []

        docs: List[Document] = []
        for slide_idx, slide in enumerate(pres.slides, start=1):
            parts: List[str] = []

            # native text
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    parts.append(shape.text)

            # speaker notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                note = slide.notes_slide.notes_text_frame.text
                if note:
                    parts.append(note)

            # OCR pictures
            if self.ocr_images and self.reader:
                for shape in slide.shapes:
                    if shape.shape_type == 13:  # PICTURE
                        try:
                            img = Image.open(io.BytesIO(shape.image.blob)).convert("RGB")
                            txt_blocks = self.reader.readtext(img, detail=0, paragraph=True)
                            if txt_blocks:
                                parts.append("\n".join(txt_blocks))
                        except Exception as e:  # noqa: BLE001
                            logger.warning("OCR failed on slide %d image: %s", slide_idx, e)

            slide_text = "\n".join(p for p in parts if p.strip())
            if not slide_text:
                continue

            # Split if requested
            chunk_texts = self.splitter.split_text(slide_text) if self.splitter else [slide_text]

            for chunk in chunk_texts:
                docs.append(
                    Document(
                        page_content=chunk,
                        metadata={
                            "source": file_path,
                            "slide": slide_idx,
                            "type": "pptx",
                        },
                    )
                )

        logger.info("Loaded %d document chunk(s) from %s", len(docs), file_path)
        return docs
